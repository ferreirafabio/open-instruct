"""Convert a meeting markdown file (this dir's format) to a .pptx deck.

Slide schema we expect inside the markdown:
  # <Deck title>            -> title slide
  optional **Presenter:** / **Slides:** / reference list before first '---'
  ## Slide N — <Title>      -> one content slide per occurrence
    - bullet
    - bullet (supports **bold** -> bold runs)
    ![alt](path)            -> image; if a slide has bullets AND an image,
                               bullets go on the left, image on the right.
                               Path is resolved relative to the markdown file.
  ## <Other heading>        -> also rendered as a content slide
  --- separators are ignored

Usage:
    uv run --with python-pptx python md_to_pptx.py <input.md> <output.pptx>
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
LINK_RE = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")
IMG_RE = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")


def parse_md(md_path: Path):
    text = md_path.read_text(encoding="utf-8")
    lines = text.splitlines()

    deck_title = ""
    intro_lines: list[str] = []
    slides: list[dict] = []
    cur: dict | None = None
    seen_first_section = False

    def flush():
        nonlocal cur
        if cur is not None:
            slides.append(cur)
        cur = None

    for raw in lines:
        line = raw.rstrip()
        if not deck_title and line.startswith("# "):
            deck_title = line[2:].strip()
            continue
        if line.startswith("## "):
            flush()
            cur = {"title": line[3:].strip(), "bullets": [], "images": []}
            seen_first_section = True
            continue
        if line.strip() == "---":
            continue
        if cur is None:
            if seen_first_section:
                continue
            if line.strip():
                intro_lines.append(line)
            continue
        # Image line
        img_m = IMG_RE.search(line)
        if img_m:
            img_path = img_m.group(2).strip()
            cur["images"].append({"alt": img_m.group(1), "path": img_path})
            continue
        m = re.match(r"\s*[-*]\s+(.*)", line)
        if m:
            cur["bullets"].append(m.group(1).strip())
            continue
        # Continuation line: append to last bullet if present
        if line.strip() and cur["bullets"]:
            cur["bullets"][-1] += " " + line.strip()
    flush()
    return deck_title, intro_lines, slides


def add_runs(paragraph, text: str, bullet_size_pt: int = 18):
    text = LINK_RE.sub(r"\1", text)
    parts = []
    last = 0
    for m in BOLD_RE.finditer(text):
        if m.start() > last:
            parts.append((text[last : m.start()], False))
        parts.append((m.group(1), True))
        last = m.end()
    if last < len(text):
        parts.append((text[last:], False))
    if not parts:
        parts = [(text, False)]

    first = True
    for chunk, bold in parts:
        if first and paragraph.runs:
            run = paragraph.runs[0]
        else:
            run = paragraph.add_run()
        run.text = chunk
        run.font.bold = bold
        run.font.size = Pt(bullet_size_pt)
        first = False


def add_bullets_textbox(slide, bullets, left, top, width, height, size_pt=18):
    from pptx.util import Emu

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.text = ""
        add_runs(p, "• " + b, bullet_size_pt=size_pt)
    return tb


def add_fitted_image(slide, img_path, left, top, max_w, max_h):
    from PIL import Image

    with Image.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    box_aspect = max_w / max_h
    if aspect > box_aspect:
        # constrain by width
        w = max_w
        h = int(max_w / aspect)
    else:
        h = max_h
        w = int(max_h * aspect)
    # center within box
    cx = left + (max_w - w) // 2
    cy = top + (max_h - h) // 2
    return slide.shapes.add_picture(str(img_path), cx, cy, width=w, height=h)


def build_pptx(deck_title: str, intro_lines: list[str], slides, out_path: Path, md_dir: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = deck_title or out_path.stem
    if len(s.placeholders) > 1:
        sub = s.placeholders[1]
        sub.text = "\n".join(l for l in intro_lines if l.strip())

    blank_layout = prs.slide_layouts[6]  # blank
    title_only_layout = prs.slide_layouts[5]  # title only

    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height
    MARGIN = Inches(0.5)
    TITLE_H = Inches(1.0)

    for spec in slides:
        title = spec["title"]
        bullets = spec["bullets"]
        images = spec["images"]

        # Resolve image paths
        resolved = []
        for img in images:
            p = (md_dir / img["path"]).resolve()
            if p.exists():
                resolved.append(p)
            else:
                print(f"WARN: image not found, skipping: {p}", file=sys.stderr)

        s = prs.slides.add_slide(title_only_layout)
        s.shapes.title.text = title

        body_top = TITLE_H + Inches(0.2)
        body_h = SLIDE_H - body_top - MARGIN

        if bullets and resolved:
            # split: bullets left, image right
            half_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) // 2
            add_bullets_textbox(
                s,
                bullets,
                left=MARGIN,
                top=body_top,
                width=half_w,
                height=body_h,
                size_pt=16,
            )
            add_fitted_image(
                s,
                resolved[0],
                left=MARGIN + half_w + Inches(0.3),
                top=body_top,
                max_w=half_w,
                max_h=body_h,
            )
        elif bullets:
            add_bullets_textbox(
                s,
                bullets,
                left=MARGIN,
                top=body_top,
                width=SLIDE_W - 2 * MARGIN,
                height=body_h,
                size_pt=18,
            )
        elif resolved:
            add_fitted_image(
                s,
                resolved[0],
                left=MARGIN,
                top=body_top,
                max_w=SLIDE_W - 2 * MARGIN,
                max_h=body_h,
            )

    prs.save(out_path)


def main():
    if len(sys.argv) != 3:
        print(__doc__)
        sys.exit(2)
    in_md = Path(sys.argv[1]).resolve()
    out_pptx = Path(sys.argv[2]).resolve()
    deck_title, intro, slides = parse_md(in_md)
    build_pptx(deck_title, intro, slides, out_pptx, in_md.parent)
    print(f"Wrote {out_pptx} ({len(slides)} content slides)")


if __name__ == "__main__":
    main()
