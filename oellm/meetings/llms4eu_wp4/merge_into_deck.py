"""Merge our generated content slides into an existing Google-Slides export.

Usage:
    uv run --no-project --with python-pptx --with pillow python merge_into_deck.py \
        <template.pptx> <our_md.md> <output.pptx>

Behaviour:
- Opens the template .pptx (e.g. the FBK deck downloaded as pptx).
- For each `## Slide N — <Title>` section in our_md.md, appends a new slide
  using the template's BLANK layout, then draws the title and body manually
  so we don't fight the template's placeholders.
- Detects the existing References slide in the template (title contains
  "Reference") and moves it to the end after insertion, so order becomes:
    cover, intro, [our slides], references.
- Each content slide gets bullets on the left and (optional) image on the
  right; everything is sized to the *actual* slide dimensions of the
  template (so it works with both 10x5.625 and 13.33x7.5 decks).
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
LINK_RE = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")
IMG_RE = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")


def parse_md(md_path: Path):
    text = md_path.read_text(encoding="utf-8")
    lines = text.splitlines()

    slides: list[dict] = []
    cur: dict | None = None

    def flush():
        nonlocal cur
        if cur is not None:
            slides.append(cur)
        cur = None

    for raw in lines:
        line = raw.rstrip()
        if line.startswith("## "):
            flush()
            cur = {"title": line[3:].strip(), "bullets": [], "images": []}
            continue
        if cur is None:
            continue
        if line.strip() == "---":
            continue
        img_m = IMG_RE.search(line)
        if img_m:
            cur["images"].append({"alt": img_m.group(1), "path": img_m.group(2).strip()})
            continue
        m = re.match(r"\s*[-*]\s+(.*)", line)
        if m:
            cur["bullets"].append(m.group(1).strip())
            continue
        if line.strip() and cur["bullets"]:
            cur["bullets"][-1] += " " + line.strip()
    flush()
    return slides


def _split_for_links(text: str):
    """Split text into [(chunk, link_or_None)] preserving markdown link labels + URLs.
    Bare http(s)://... URLs are also treated as links (label == url).
    """
    out = []
    pos = 0
    pattern = re.compile(r"\[([^\]]+)\]\(([^)]+)\)|(https?://\S+)")
    for m in pattern.finditer(text):
        if m.start() > pos:
            out.append((text[pos : m.start()], None))
        if m.group(1):
            out.append((m.group(1), m.group(2)))
        else:
            url = m.group(3).rstrip(".,;:)]")
            out.append((url, url))
        pos = m.end()
    if pos < len(text):
        out.append((text[pos:], None))
    if not out:
        out = [(text, None)]
    return out


def _split_for_bold(text: str):
    """Split a (text, link) chunk further by **bold** markers."""
    parts = []
    last = 0
    for m in BOLD_RE.finditer(text):
        if m.start() > last:
            parts.append((text[last : m.start()], False))
        parts.append((m.group(1), True))
        last = m.end()
    if last < len(text):
        parts.append((text[last:], False))
    if not parts:
        parts = [(text, False)]
    return parts


def add_runs(paragraph, text: str, size_pt: int, bold_default: bool = False, color: RGBColor | None = None):
    link_parts = _split_for_links(text)
    first = True
    for chunk, link in link_parts:
        for sub_chunk, bold in _split_for_bold(chunk):
            if first and paragraph.runs:
                run = paragraph.runs[0]
            else:
                run = paragraph.add_run()
            run.text = sub_chunk
            run.font.bold = bold or bold_default
            run.font.size = Pt(size_pt)
            if link is not None:
                run.hyperlink.address = link
            elif color is not None:
                run.font.color.rgb = color
            first = False


def add_title_textbox(slide, text: str, left, top, width, height, size_pt: int, color: RGBColor | None = None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.text = ""
    add_runs(p, text, size_pt=size_pt, bold_default=True, color=color)
    return tb


def add_bullets_textbox(slide, bullets, left, top, width, height, size_pt):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.text = ""
        add_runs(p, "• " + b, size_pt=size_pt)
        p.space_after = Pt(4)
    return tb


def add_fitted_image(slide, img_path, left, top, max_w, max_h):
    from PIL import Image

    with Image.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    box_aspect = max_w / max_h
    if aspect > box_aspect:
        w = max_w
        h = int(max_w / aspect)
    else:
        h = max_h
        w = int(max_h * aspect)
    cx = left + (max_w - w) // 2
    cy = top + (max_h - h) // 2
    return slide.shapes.add_picture(str(img_path), cx, cy, width=w, height=h)


def pick_blank_layout(prs):
    for lay in prs.slide_layouts:
        if (lay.name or "").upper() == "BLANK":
            return lay
    # fallback: layout with the fewest placeholders
    return min(prs.slide_layouts, key=lambda l: len(list(l.placeholders)))


def get_slide_idx_by_title_substring(prs, substr: str) -> int | None:
    needle = substr.lower()
    for i, sl in enumerate(prs.slides):
        if sl.shapes.title and sl.shapes.title.text and needle in sl.shapes.title.text.lower():
            return i
    return None


def move_slide(prs, old_idx: int, new_idx: int):
    sld_id_lst = prs.slides._sldIdLst
    slides = list(sld_id_lst)
    elem = slides[old_idx]
    sld_id_lst.remove(elem)
    sld_id_lst.insert(new_idx, elem)


def main():
    if len(sys.argv) != 4:
        print(__doc__)
        sys.exit(2)
    template_pptx = Path(sys.argv[1]).resolve()
    md_path = Path(sys.argv[2]).resolve()
    out_pptx = Path(sys.argv[3]).resolve()
    md_dir = md_path.parent

    prs = Presentation(str(template_pptx))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    # Inches at the actual template scale
    inch_w = slide_w / 914400
    inch_h = slide_h / 914400
    print(f"Template: {len(prs.slides)} slides, {inch_w:.2f}x{inch_h:.2f} in.")

    layout = pick_blank_layout(prs)
    print(f"Using layout: {layout.name}")

    # Layout constants — sized to actual slide dimensions
    margin = Emu(int(0.4 * 914400))  # 0.4 in
    title_top = Emu(int(0.25 * 914400))
    title_h = Emu(int(0.55 * 914400))
    title_size_pt = 22 if inch_w < 11 else 28
    bullet_size_pt = 13 if inch_w < 11 else 16

    body_top = title_top + title_h + Emu(int(0.15 * 914400))
    body_left = margin
    body_w = slide_w - 2 * margin
    body_h = slide_h - body_top - margin

    # FBK-ish dark blue title color (close enough to the deck's accent)
    title_color = RGBColor(0x1A, 0x3A, 0x6E)

    md_slides = parse_md(md_path)
    # Drop the template's existing References slide; we'll add our own at the end
    refs_idx = get_slide_idx_by_title_substring(prs, "reference")
    if refs_idx is not None:
        sld_id_lst = prs.slides._sldIdLst
        elem = list(sld_id_lst)[refs_idx]
        rId = elem.rId
        prs.part.drop_rel(rId)
        sld_id_lst.remove(elem)
        print(f"Removed template references slide at idx {refs_idx}.")
    print(f"Inserting {len(md_slides)} content slides.")

    for spec in md_slides:
        title = spec["title"]
        bullets = spec["bullets"]
        images = spec["images"]

        resolved = []
        for img in images:
            p = (md_dir / img["path"]).resolve()
            if p.exists():
                resolved.append(p)
            else:
                print(f"WARN: image not found, skipping: {p}", file=sys.stderr)

        s = prs.slides.add_slide(layout)

        add_title_textbox(
            s, title,
            left=margin, top=title_top, width=slide_w - 2 * margin, height=title_h,
            size_pt=title_size_pt, color=title_color,
        )

        if bullets and resolved:
            gap = Emu(int(0.2 * 914400))
            half_w = (body_w - gap) // 2
            add_bullets_textbox(
                s, bullets,
                left=body_left, top=body_top, width=half_w, height=body_h,
                size_pt=bullet_size_pt,
            )
            add_fitted_image(
                s, resolved[0],
                left=body_left + half_w + gap, top=body_top,
                max_w=half_w, max_h=body_h,
            )
        elif bullets:
            add_bullets_textbox(
                s, bullets,
                left=body_left, top=body_top, width=body_w, height=body_h,
                size_pt=bullet_size_pt + 2,
            )
        elif resolved:
            add_fitted_image(
                s, resolved[0],
                left=body_left, top=body_top, max_w=body_w, max_h=body_h,
            )

    prs.save(str(out_pptx))
    print(f"Saved {out_pptx} with {len(prs.slides)} slides total.")


if __name__ == "__main__":
    main()
